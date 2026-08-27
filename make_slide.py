#!/usr/bin/env python3
"""สร้าง workshop_slide.pptx จากเนื้อหาเดียวกับ workshop_slide.html (6 หน้า)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ธีมสี
BG      = RGBColor(0x0F, 0x11, 0x17)
BLUE    = RGBColor(0x7C, 0x9C, 0xFF)
PURPLE  = RGBColor(0xA7, 0x8B, 0xFA)
GREEN   = RGBColor(0x34, 0xD3, 0x99)
BLUEB   = RGBColor(0x60, 0xA5, 0xFA)
PURPLEB = RGBColor(0xC0, 0x84, 0xFC)
WHITE   = RGBColor(0xE8, 0xEA, 0xED)
GREY    = RGBColor(0x9D, 0xB4, 0xFF)
DARK    = RGBColor(0x1A, 0x1F, 0x2E)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def add_slide():
    s = prs.slides.add_slide(blank)
    bg = s.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return s

def textbox(s, l, t, w, h, lines):
    """lines = list of (text, size, color, bold)"""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (txt, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Leelawadee UI"
        p.space_after = Pt(8)
    return tb

def bullet(s, l, t, w, h, items, size=28, color=WHITE):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + it
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Leelawadee UI"
        p.space_after = Pt(10)
    return tb

# ---- Slide 1: ปก ----
s = add_slide()
textbox(s, 0.8, 2.4, 11.7, 2.0, [
    ("Hermes Agent กับ 3 บอทส่วนตัว", 54, BLUE, True),
])
textbox(s, 0.8, 4.2, 11.7, 1.0, [
    ("อบรม AI Agent สำหรับนิสิต ICT · 90 นาที", 30, WHITE, False),
])
# tags
tags = ["ไม่ต้องเขียนโค้ด", "พิมพ์คุยได้งาน", "ใช้งานจริงได้ทันที"]
x = 0.8
for tg in tags:
    tb = s.shapes.add_textbox(Inches(x), Inches(5.4), 3.4, 0.6)
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = tg
    p.font.size = Pt(18); p.font.color.rgb = GREY; p.font.name = "Leelawadee UI"
    x += 3.6

# ---- Slide 2: AI Agent vs ChatGPT ----
s = add_slide()
textbox(s, 0.8, 0.6, 11.7, 1.0, [("AI Agent ต่างจาก ChatGPT ยังไง?", 40, PURPLE, True)])
bullet(s, 0.8, 2.0, 11.7, 4.5, [
    "ChatGPT → แค่ตอบข้อความ",
    "Hermes Agent → มีมือทำงานได้จริง: อ่านไฟล์, สร้างสไลด์, วิเคราะห์ข้อมูล, วาดรูป",
    "ตั้งเป็นบอทเฉพาะงานได้ → แต่ละตัวมีบุคลิกและหน้าที่ของตัวเอง",
], size=30)

# ---- Slide 3: 3 bots ----
s = add_slide()
textbox(s, 0.8, 0.5, 11.7, 1.0, [("วันนี้เรามี 3 บอทให้นิสิตลองคุย", 38, PURPLE, True)])
cards = [
    (GREEN,   "🟢 พี่ติว", "สรุปบทเรียน · ทำข้อสอบ · ตรวจงาน", "study-coach"),
    (BLUEB,   "🔵 นักวิเคราะห์", "วิเคราะห์ข้อมูล · ทำกราฟ · หา insight", "data-analyst"),
    (PURPLEB, "🟣 นักเขียน", "ทำสไลด์ · วาดรูปประกอบ", "content-creator"),
]
x = 0.8
for color, title, desc, cmd in cards:
    box = s.shapes.add_shape(1, Inches(x), Inches(2.0), Inches(3.7), Inches(3.4))
    box.fill.solid(); box.fill.fore_color.rgb = DARK
    box.line.color.rgb = color; box.line.width = Pt(3)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = color; p.font.name = "Leelawadee UI"
    p2 = tf.add_paragraph(); p2.text = desc
    p2.font.size = Pt(20); p2.font.color.rgb = WHITE; p2.font.name = "Leelawadee UI"; p2.space_before = Pt(10)
    p3 = tf.add_paragraph(); p3.text = cmd
    p3.font.size = Pt(16); p3.font.color.rgb = GREEN; p3.font.name = "Consolas"; p3.space_before = Pt(14)
    x += 3.95

# ---- Slide 4: โครงเวลา + Demo ----
s = add_slide()
textbox(s, 0.8, 0.5, 11.7, 1.0, [("รูปแบบเวิร์กชอป (90 นาที)", 38, PURPLE, True)])
bullet(s, 0.8, 2.0, 11.7, 3.0, [
    "0–15 นาที: Intro + อธิบายแนวคิด",
    "15–30 นาที: นิสิตตั้งบอท 3 ตัวพร้อมกัน",
    "30–84 นาที: ทุกคนลองคุยกับบอททีละตัว (พี่ติว → นักวิเคราะห์ → นักเขียน)",
    "84–90 นาที: Showcase + Q&A",
], size=26)
box = s.shapes.add_shape(1, Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.2))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x24,0x1A,0x2E)
box.line.color.rgb = PURPLEB; box.line.width = Pt(2)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Demo พิเศษ (วิทยากร): ปฏิทิน + อีเมลอัตโนมัติ (ต้อง login Google → ห้ามทำบนเครื่องรวม)"
p.font.size = Pt(20); p.font.color.rgb = WHITE; p.font.name = "Leelawadee UI"

# ---- Slide 5: กฎทอง ----
s = add_slide()
textbox(s, 0.8, 0.5, 11.7, 1.0, [("กฎทองใช้ AI Agent", 40, PURPLE, True)])
bullet(s, 0.8, 2.0, 11.7, 4.5, [
    "1. บอกบริบท — ใคร ทำไม ให้ใครดู",
    "2. ขอรูปแบบ — สรุป bullet / ตาราง / สไลด์",
    "3. ตรวจเสมอ — AI อาจผิด อย่าเชื่อ 100%",
    "4. ห้ามใส่รหัส/ข้อมูลส่วนตัว",
], size=30)

# ---- Slide 6: จบ ----
s = add_slide()
textbox(s, 0.8, 2.8, 11.7, 1.5, [
    ("พร้อมลุยกันเลย! 🚀", 52, BLUE, True),
])
textbox(s, 0.8, 4.4, 11.7, 1.0, [
    ("เปิดแท็บ Bots ใน Hermes Desktop แล้วคุยกับบอททั้ง 3 ตัวได้เลย", 28, WHITE, False),
])

out = "workshop_slide.pptx"
prs.save(out)
print("สร้างไฟล์:", out, "จำนวนหน้า:", len(prs.slides._sldIdLst))
